from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
OUT = DOCS / "reports" / "compact_training_summary" / "强化学习训练流程_三页精简汇报.pptx"
VISUAL = DOCS / "visualizations" / "hover_policy_scene_short" / "policy_scene_screenshot.png"
SUMMARY = ROOT / "experiments" / "hover_rl_reproduction" / "results" / "repro_hover_short" / "summary.json"

NAVY = RGBColor(20, 35, 54)
BLUE = RGBColor(35, 93, 164)
CYAN = RGBColor(82, 173, 198)
GREEN = RGBColor(57, 145, 112)
AMBER = RGBColor(216, 146, 45)
RED = RGBColor(190, 77, 72)
INK = RGBColor(32, 41, 52)
MUTED = RGBColor(92, 108, 124)
LINE = RGBColor(202, 215, 230)
PANEL = RGBColor(244, 248, 252)
WHITE = RGBColor(255, 255, 255)


def summary() -> dict:
    if SUMMARY.exists():
        return json.loads(SUMMARY.read_text(encoding="utf-8"))
    return {
        "timesteps": 2048,
        "mean_reward": 355.7683,
        "std_reward": 0.1818,
        "final_z": 0.2163,
    }


def blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def text(slide, x, y, w, h, value, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def title(slide, value, kicker="gym-pybullet-drones | PPO hover control"):
    text(slide, 0.85, 0.25, 7.0, 0.35, kicker, size=8.5, color=BLUE, bold=True)
    text(slide, 0.85, 0.78, 18.5, 0.95, value, size=21, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.85), Cm(1.83), Cm(23.7), Cm(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def pill(slide, x, y, w, h, label, fill, color=WHITE, size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return shape


def card(slide, x, y, w, h, heading, body, accent=BLUE, body_size=11):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = PANEL
    s.line.color.rgb = RGBColor(218, 228, 239)
    s.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(0.08), Cm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text(slide, x + 0.32, y + 0.22, w - 0.55, 0.35, heading, size=12, color=NAVY, bold=True)
    text(slide, x + 0.32, y + 0.72, w - 0.55, h - 0.85, body, size=body_size, color=INK)
    return s


def arrow(slide, x1, y1, x2, y2, color=BLUE):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.color.rgb = color
    c.line.width = Pt(2)
    c.line.end_arrowhead = True
    return c


def stage(slide, x, y, w, h, num, name, desc, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    pill(slide, x + 0.18, y + 0.18, 0.62, 0.62, str(num), fill, size=12)
    text(slide, x + 0.92, y + 0.20, w - 1.05, 0.35, name, size=10.2, color=NAVY, bold=True)
    text(slide, x + 0.28, y + 0.78, w - 0.55, h - 0.9, desc, size=8.0, color=MUTED)
    return shape


def metric(slide, x, y, label, value, note, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(4.75), Cm(2.05))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = LINE
    text(slide, x + 0.25, y + 0.20, 4.2, 0.3, label, size=9.5, color=MUTED, bold=True)
    text(slide, x + 0.25, y + 0.62, 4.2, 0.55, value, size=18, color=color, bold=True)
    text(slide, x + 0.25, y + 1.35, 4.2, 0.45, note, size=8.0, color=MUTED)


def add_slide_1(prs):
    slide = blank(prs)
    title(slide, "这个项目到底在训练什么？")
    text(slide, 0.9, 2.12, 9.4, 0.65, "用奖励信号训练策略，而不是手写控制律。", size=18, color=INK, bold=True)
    if VISUAL.exists():
        slide.shapes.add_picture(str(VISUAL), Cm(11.0), Cm(2.03), width=Cm(12.9), height=Cm(7.25))
    card(slide, 0.9, 3.05, 4.6, 2.0, "目标", "悬停到\nx=0, y=0, z=1", BLUE, 12.5)
    card(slide, 5.85, 3.05, 4.6, 2.0, "策略", "PPO 网络\n状态 -> 动作", GREEN, 12.5)
    card(slide, 0.9, 5.50, 4.6, 2.0, "环境", "PyBullet\n四旋翼仿真", CYAN, 12.5)
    card(slide, 5.85, 5.50, 4.6, 2.0, "奖励", "接近目标\n分数升高", AMBER, 12.5)
    text(slide, 0.95, 9.55, 22.8, 0.55, "本质：不是手写控制律，而是用奖励信号驱动策略从试错中学习。", size=17, color=NAVY, bold=True)
    pill(slide, 0.9, 10.55, 4.5, 0.7, "obs 读取状态", BLUE)
    arrow(slide, 5.55, 10.9, 6.75, 10.9)
    pill(slide, 6.95, 10.55, 4.5, 0.7, "action 输出动作", GREEN)
    arrow(slide, 11.6, 10.9, 12.8, 10.9)
    pill(slide, 13.0, 10.55, 4.5, 0.7, "reward 评价", AMBER)
    arrow(slide, 17.65, 10.9, 18.85, 10.9)
    pill(slide, 19.05, 10.55, 4.5, 0.7, "learn 更新策略", RED)


def add_slide_2(prs):
    slide = blank(prs)
    title(slide, "训练链路：从环境到模型，再到展示")
    y = 2.85
    xs = [0.85, 4.75, 8.65, 12.55, 16.45, 20.35]
    stages = [
        ("Docker", "固定环境", BLUE),
        ("Hover", "任务/奖励", CYAN),
        ("RL接口", "obs/action", GREEN),
        ("PPO", "更新策略", AMBER),
        ("评估", "保存模型", RED),
        ("展示", "截图/视频", BLUE),
    ]
    for i, (name, desc, color) in enumerate(stages, 1):
        stage(slide, xs[i - 1], y, 3.25, 2.15, i, name, desc, color)
        if i < len(stages):
            arrow(slide, xs[i - 1] + 3.32, y + 1.08, xs[i] - 0.15, y + 1.08)
    card(slide, 0.95, 6.05, 6.9, 2.18, "快速复现", "experiments/hover_rl_reproduction/scripts/reproduce_hover_short.py\n训练 / 评估 / 导出结果", BLUE, 10.2)
    card(slide, 8.25, 6.05, 6.9, 2.18, "正式训练", "examples/learn.py\n周期评估 / 保存 best_model", GREEN, 10.2)
    card(slide, 15.55, 6.05, 7.7, 2.18, "策略展示", "play.py 或 render_policy_scene.py\nGUI / 截图 / MP4", AMBER, 10.2)
    text(slide, 1.0, 9.35, 22.2, 0.9, "训练闭环 = 任务定义 + RL 接口 + PPO 算法 + 评估保存 + 可视化展示。", size=18, color=NAVY, bold=True)


def add_slide_3(prs):
    s = summary()
    slide = blank(prs)
    title(slide, "怎么跑、怎么调、怎么算达标")
    card(slide, 0.85, 2.10, 7.0, 2.65, "部署", "启动 Docker\nbuild 镜像\nrun 训练脚本", BLUE, 11)
    card(slide, 8.3, 2.10, 7.0, 2.65, "核心参数", "timesteps 训练步数\neval-episodes 评估回合\nrollout-steps 轨迹长度", GREEN, 11)
    card(slide, 15.75, 2.10, 7.8, 2.65, "验收逻辑", "奖励 + 高度误差 + 轨迹稳定\n+ 截断状态 + 视频表现", AMBER, 11)
    metric(slide, 1.0, 5.90, "当前短训练", f"{int(s.get('timesteps', 2048))}", "steps，用于验证链路", BLUE)
    metric(slide, 6.1, 6.2, "mean_reward", f"{float(s.get('mean_reward', 0)):.1f}", "目标约 474", GREEN)
    metric(slide, 11.2, 6.2, "final_z", f"{float(s.get('final_z', 0)):.3f}", "目标 z=1.0", AMBER)
    metric(slide, 16.3, 6.2, "状态判断", "未达标", "已跑通，需长训练", RED)
    text(slide, 1.0, 9.2, 5.0, 0.4, "建议达标线", size=13, color=NAVY, bold=True)
    pill(slide, 1.0, 9.8, 4.8, 0.72, "reward >= 450", GREEN, size=10.5)
    pill(slide, 6.2, 9.8, 4.8, 0.72, "|z - 1| <= 0.1", BLUE, size=10.5)
    pill(slide, 11.4, 9.8, 5.2, 0.72, "成功率 >= 80%", AMBER, size=10.5)
    pill(slide, 17.0, 9.8, 5.8, 0.72, "不触发失控截断", RED, size=10.5)
    text(slide, 1.0, 11.18, 22.0, 0.55, "结论：当前素材适合展示“流程已跑通”；若要展示“控制器性能好”，下一步应跑 1e6+ 步并做多随机种子评估。", size=16, color=NAVY, bold=True)


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.288)
    add_slide_1(prs)
    add_slide_2(prs)
    add_slide_3(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
