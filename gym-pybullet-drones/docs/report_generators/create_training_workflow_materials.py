from __future__ import annotations

import json
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib import font_manager
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm as PptCm
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets" / "training_workflow"
VISUAL = DOCS / "visualizations" / "hover_policy_scene_short" / "policy_scene_screenshot.png"
WORD_OUT = DOCS / "reports" / "training_workflow" / "强化学习训练工作流程说明.docx"
PPT_OUT = DOCS / "reports" / "training_workflow" / "强化学习训练工作流程汇报.pptx"
REPRO_RESULTS = ROOT / "experiments" / "hover_rl_reproduction" / "results"


TITLE_COLOR = RGBColor(27, 62, 100)
ACCENT = RGBColor(54, 113, 178)
LIGHT = RGBColor(230, 240, 250)
TEXT = RGBColor(28, 35, 43)

CHINESE_FONT = Path("C:/Windows/Fonts/msyh.ttc")
FONT_PROPS = font_manager.FontProperties(fname=str(CHINESE_FONT)) if CHINESE_FONT.exists() else None


def ensure_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)


def set_run_font(run, size: float | None = None, bold: bool = False) -> None:
    run.font.name = "Microsoft YaHei"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    if size:
        run.font.size = Pt(size)
    run.bold = bold


def setup_doc(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Cm(2.0)
    section.bottom_margin = Cm(1.8)
    section.left_margin = Cm(2.2)
    section.right_margin = Cm(2.2)
    normal = document.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)
    for name in ["Heading 1", "Heading 2", "Heading 3"]:
        style = document.styles[name]
        style.font.name = "Microsoft YaHei"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")


def add_para(document: Document, text: str = "", bold: bool = False):
    p = document.add_paragraph()
    p.paragraph_format.line_spacing = 1.15
    p.paragraph_format.space_after = Pt(6)
    if text:
        run = p.add_run(text)
        set_run_font(run, bold=bold)
    return p


def add_bullets(document: Document, items: list[str]) -> None:
    for item in items:
        p = document.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        run = p.add_run(item)
        set_run_font(run)


def add_numbered(document: Document, items: list[str]) -> None:
    for item in items:
        p = document.add_paragraph(style="List Number")
        p.paragraph_format.space_after = Pt(3)
        run = p.add_run(item)
        set_run_font(run)


def add_code(document: Document, text: str) -> None:
    p = document.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.5)
    p.paragraph_format.space_after = Pt(8)
    run = p.add_run(text)
    run.font.name = "Consolas"
    run.font.size = Pt(9)


def shade_cell(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell(cell, text: str, bold: bool = False) -> None:
    cell.text = ""
    run = cell.paragraphs[0].add_run(text)
    set_run_font(run, bold=bold)


def add_picture(document: Document, path: Path, width_cm: float) -> None:
    if not path.exists():
        return
    p = document.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(str(path), width=Cm(width_cm))
    p.paragraph_format.space_after = Pt(8)


def render_formula(source: str, name: str, width: float = 6.5) -> Path:
    path = ASSETS / name
    fig = plt.figure(figsize=(width, 0.8), dpi=220)
    fig.patch.set_alpha(0)
    plt.text(0.5, 0.5, source, ha="center", va="center", fontsize=18)
    plt.axis("off")
    plt.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    return path


def render_workflow_chart() -> Path:
    path = ASSETS / "rl_workflow.png"
    fig, ax = plt.subplots(figsize=(11, 4.4), dpi=180)
    ax.axis("off")
    labels = [
        "reset\n创建环境",
        "obs\n读取状态",
        "policy\n输出动作",
        "step\n仿真推进",
        "reward\n评价好坏",
        "learn\n更新策略",
        "save/eval\n保存评估",
    ]
    xs = [0.08, 0.22, 0.36, 0.50, 0.64, 0.78, 0.92]
    y = 0.55
    for i, (x, label) in enumerate(zip(xs, labels)):
        ax.text(
            x,
            y,
            label,
            ha="center",
            va="center",
            fontsize=11,
            color="#1b3e64",
            fontproperties=FONT_PROPS,
            bbox=dict(boxstyle="round,pad=0.45", fc="#e6f0fa", ec="#3671b2", lw=1.6),
        )
        if i < len(xs) - 1:
            ax.annotate(
                "",
                xy=(xs[i + 1] - 0.055, y),
                xytext=(x + 0.055, y),
                arrowprops=dict(arrowstyle="->", color="#3671b2", lw=1.8),
            )
    ax.annotate(
        "大量重复试飞，策略逐渐变好",
        xy=(0.50, 0.28),
        xytext=(0.50, 0.14),
        ha="center",
        fontsize=12,
        color="#2d5f91",
        fontproperties=FONT_PROPS,
        arrowprops=dict(arrowstyle="->", color="#2d5f91", lw=1.6),
    )
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def render_role_chart() -> Path:
    path = ASSETS / "software_roles.png"
    fig, ax = plt.subplots(figsize=(10, 5.6), dpi=180)
    ax.axis("off")
    rows = [
        ("Docker", "固定 Linux + Python 3.10 环境，保证复现稳定"),
        ("PyBullet", "提供四旋翼物理仿真和 3D 场景"),
        ("Gymnasium", "把仿真封装成 reset/step/render 标准接口"),
        ("Stable-Baselines3", "提供 PPO 算法，实现训练、评估和模型保存"),
        ("PyTorch", "作为策略神经网络的底层计算框架"),
        ("项目脚本", "定义任务、奖励、参数、训练入口和可视化输出"),
    ]
    y = 0.90
    for name, desc in rows:
        ax.text(0.08, y, name, fontsize=13, weight="bold", color="#1b3e64", va="center", fontproperties=FONT_PROPS)
        ax.text(0.27, y, desc, fontsize=12, color="#1c232b", va="center", fontproperties=FONT_PROPS)
        ax.plot([0.06, 0.94], [y - 0.055, y - 0.055], color="#d0d9e6", lw=1)
        y -= 0.14
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def load_summary() -> dict:
    path = REPRO_RESULTS / "repro_hover_short" / "summary.json"
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def build_word() -> None:
    ensure_assets()
    WORD_OUT.parent.mkdir(parents=True, exist_ok=True)
    workflow = render_workflow_chart()
    roles = render_role_chart()
    target_formula = render_formula(r"$p^\ast=[0,\ 0,\ 1]^\mathsf{T}$", "target_formula.png", 4.5)
    reward_formula = render_formula(
        r"$r_t=\max\left(0,\ 2-\left\|p^\ast-p_t\right\|_2^4\right)$",
        "reward_formula.png",
        7.0,
    )
    rpm_formula = render_formula(
        r"$\mathrm{RPM}_i=\mathrm{RPM}_{hover}\left(1+0.05a_t\right)$",
        "rpm_formula.png",
        6.2,
    )
    summary = load_summary()

    document = Document()
    setup_doc(document)
    title = document.add_heading("强化学习训练工作流程说明", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_para(document, "项目对象：gym-pybullet-drones 四旋翼无人机 PPO 悬停控制复现与展示。")

    document.add_heading("1. 先理解强化学习在做什么", level=1)
    add_para(document, "强化学习不是把控制公式直接写死，而是让智能体在环境里反复试错。每一次试飞，环境都会告诉它当前状态、执行动作后的结果以及奖励分数。算法根据奖励调整策略，使下一次决策更可能得到高分。")
    add_bullets(document, [
        "智能体：PPO 策略网络，也就是训练后的无人机控制策略。",
        "环境：PyBullet 中的四旋翼无人机仿真世界。",
        "状态：位置、姿态、速度、角速度和历史动作等观测信息。",
        "动作：策略输出的控制量，最终会被转换成电机转速。",
        "奖励：根据无人机是否接近目标悬停点给出的分数。",
        "目标：让无人机稳定接近并悬停在目标位置。",
    ])
    add_picture(document, workflow, 15.5)

    document.add_heading("2. 本项目训练的具体任务", level=1)
    add_para(document, "当前复现任务是单架四旋翼无人机悬停控制。目标位置定义在 HoverAviary.py 中：")
    add_picture(document, target_formula, 4.2)
    add_para(document, "也就是让无人机尽量停在 x=0、y=0、z=1 的位置。")
    add_para(document, "奖励函数鼓励无人机靠近目标点：")
    add_picture(document, reward_formula, 7.5)
    add_para(document, "距离目标越近，奖励越高；距离越远，奖励越低。短训练模型可以证明链路跑通，但还不能代表最终控制效果。")
    add_picture(document, VISUAL, 15.5)

    document.add_heading("3. 哪些软件发挥作用", level=1)
    add_picture(document, roles, 15.0)
    add_para(document, "其中 PyBullet 和 Gazebo 的角色最接近，都是仿真器。但本项目选择 PyBullet，是因为它轻量、适合快速训练和强化学习循环。")

    document.add_heading("4. 关键脚本分工", level=1)
    table = document.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    set_cell(table.rows[0].cells[0], "文件", True)
    set_cell(table.rows[0].cells[1], "作用", True)
    shade_cell(table.rows[0].cells[0], "D9EAF7")
    shade_cell(table.rows[0].cells[1], "D9EAF7")
    rows = [
        ("reproducibility/docker/Dockerfile.repro", "构建稳定的 Docker 运行环境。"),
        ("pyproject.toml", "声明 Python 依赖，例如 PyBullet、Gymnasium、Stable-Baselines3。"),
        ("HoverAviary.py", "定义悬停任务、目标位置、奖励函数和结束条件。"),
        ("BaseRLAviary.py", "定义强化学习动作空间、观测空间和动作到电机转速的转换。"),
        ("experiments/hover_rl_reproduction/scripts/reproduce_hover_short.py", "最小复现脚本，负责训练、评估、保存模型和导出轨迹。"),
        ("learn.py", "官方训练脚本，支持长训练、周期评估和 best_model 保存。"),
        ("play.py", "加载训练好的模型，在 PyBullet GUI 中播放策略。"),
        ("Logger.py", "记录和绘制位置、速度、姿态、电机转速等曲线。"),
        ("render_policy_scene.py", "生成展示用截图、GIF 和 MP4。"),
    ]
    for file, role in rows:
        cells = table.add_row().cells
        set_cell(cells[0], file)
        set_cell(cells[1], role)

    document.add_heading("5. 如何部署环境", level=1)
    add_numbered(document, [
        "安装 Docker Desktop，并确认 Docker daemon 已启动。",
        "打开 PowerShell，进入项目目录。",
        "使用 reproducibility/docker/Dockerfile.repro 构建复现镜像。",
        "运行训练命令，把结果挂载到本地 results 目录。",
    ])
    add_code(document, "cd E:\\1-AI辅助工作\\科研项目\\强化学习\\gym-pybullet-drones\n\ndocker build -f reproducibility/docker/Dockerfile.repro -t gym-pybullet-drones-repro .")

    document.add_heading("6. 如何设置参数并开始训练", level=1)
    add_para(document, "当前最小复现脚本中使用运动学观测和一维转速动作：")
    add_code(document, "OBS_TYPE = ObservationType.KIN\nACT_TYPE = ActionType.ONE_D_RPM")
    add_para(document, "一维动作会被转换为四个电机共同变化的转速命令：")
    add_picture(document, rpm_formula, 6.7)
    add_para(document, "启动短训练：")
    add_code(document, 'docker run --rm `\n  -v "${PWD}\\experiments\\hover_rl_reproduction\\scripts\\reproduce_hover_short.py:/workspace/experiments/hover_rl_reproduction/scripts/reproduce_hover_short.py:ro" `\n  -v "${PWD}\\experiments\\hover_rl_reproduction\\results:/workspace/experiments/hover_rl_reproduction/results" `\n  gym-pybullet-drones-repro `\n  python experiments/hover_rl_reproduction/scripts/reproduce_hover_short.py `\n  --timesteps 2048 `\n  --eval-episodes 3 `\n  --rollout-steps 240 `\n  --output-folder experiments/hover_rl_reproduction/results/repro_hover_short')
    add_bullets(document, [
        "--timesteps：训练步数，2048 用于快速验证，1e6 以上更接近正式实验。",
        "--eval-episodes：评估回合数，越多结果越稳定。",
        "--rollout-steps：导出测试飞行轨迹的步数。",
        "--output-folder：保存模型、摘要和轨迹的目录。",
    ])

    document.add_heading("7. 如何判断训练是否达标", level=1)
    add_para(document, "不要只看单一奖励，应同时看奖励、轨迹、高度误差、姿态稳定性、动作平滑性和可视化视频。")
    add_bullets(document, [
        "mean_reward 越高通常越好，官方 one_d_rpm 单机悬停目标奖励约为 474。",
        "final_z 应尽量接近目标高度 1.0。",
        "rollout.csv 中的 x、y 应接近 0，z 应接近 1.0。",
        "训练过程中不应频繁触发 truncated，动作不应剧烈震荡。",
        "可视化视频中应能看到无人机逐步接近目标高度并保持稳定。",
    ])
    if summary:
        add_para(document, "当前短训练结果如下，说明链路已经跑通，但策略还没有达到高质量悬停：")
        add_code(
            document,
            "\n".join(
                [
                    f"timesteps: {summary.get('timesteps')}",
                    f"mean_reward: {float(summary.get('mean_reward')):.4f}",
                    f"std_reward: {float(summary.get('std_reward')):.4f}",
                    f"final_z: {float(summary.get('final_z')):.4f}",
                    "target_z: 1.0000",
                ]
            ),
        )
    add_para(document, "建议设置正式验收标准：mean_reward >= 450，|final_z - 1.0| <= 0.1，连续 10 次评估成功率 >= 80%，且不触发失控截断。")

    document.add_heading("8. 后续正式研究路线", level=1)
    add_numbered(document, [
        "先用 2048 步确认环境能跑通。",
        "增加到 100000 步，观察策略是否开始学习。",
        "增加到 1000000 步以上，记录奖励曲线和轨迹曲线。",
        "设置多个随机种子，判断结果是否稳定。",
        "与 PID 或 DSLPIDControl 做对比。",
        "修改奖励函数，加入姿态稳定、能耗和动作平滑惩罚。",
        "加入风扰、噪声、电机延迟等扰动，测试鲁棒性。",
        "最后再考虑 SITL 或实机部署。",
    ])

    document.add_heading("9. 汇报时可以这样概括", level=1)
    add_para(document, "这个项目用 PyBullet 搭建四旋翼仿真环境，用 Gymnasium 封装强化学习接口，用 Stable-Baselines3 的 PPO 算法训练无人机悬停控制策略。训练时，智能体观察无人机状态，输出电机动作，环境根据是否接近目标高度给奖励。训练完成后，通过模型文件、summary.json、rollout.csv、曲线和 3D 场景视频展示策略效果。")

    document.add_heading("10. 公式来源附录", level=1)
    add_code(document, "p* = [0, 0, 1]^T\nr_t = max(0, 2 - ||p* - p_t||_2^4)\nRPM_i = RPM_hover * (1 + 0.05 a_t)")
    document.save(WORD_OUT)


def set_ppt_font(text_frame, size=22, bold=False, color=TEXT) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = PptPt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def add_slide_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(PptCm(0.8), PptCm(0.35), PptCm(24.0), PptCm(1.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptCm(0.8), PptCm(1.45), PptCm(24.0), PptCm(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_textbox(slide, left, top, width, height, text, size=19, bold=False, fill=None):
    shape = slide.shapes.add_textbox(PptCm(left), PptCm(top), PptCm(width), PptCm(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(size)
    p.font.bold = bold
    p.font.color.rgb = TEXT
    if fill:
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptCm(left - 0.1), PptCm(top - 0.1), PptCm(width + 0.2), PptCm(height + 0.2))
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.color.rgb = RGBColor(200, 217, 235)
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.append(shape._element)
    return shape


def add_bullet_box(slide, left, top, width, height, items: list[str], size=17):
    box = slide.shapes.add_textbox(PptCm(left), PptCm(top), PptCm(width), PptCm(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = PptPt(size)
        p.font.color.rgb = TEXT
        p.space_after = PptPt(4)
    return box


def add_pill(slide, left, top, width, height, text: str, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptCm(left), PptCm(top), PptCm(width), PptCm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(15)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    return shape


def build_ppt() -> None:
    ensure_assets()
    PPT_OUT.parent.mkdir(parents=True, exist_ok=True)
    workflow = render_workflow_chart()
    roles = render_role_chart()
    target_formula = render_formula(r"$p^\ast=[0,\ 0,\ 1]^\mathsf{T}$", "ppt_target_formula.png", 4.5)
    reward_formula = render_formula(r"$r_t=\max\left(0,\ 2-\left\|p^\ast-p_t\right\|_2^4\right)$", "ppt_reward_formula.png", 7.0)
    summary = load_summary()

    prs = Presentation()
    prs.slide_width = PptCm(25.4)
    prs.slide_height = PptCm(14.288)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 1.0, 1.2, 15.0, 1.4, "强化学习训练工作流程说明", size=34, bold=True)
    add_textbox(slide, 1.05, 2.55, 22.4, 0.8, "gym-pybullet-drones 四旋翼 PPO 悬停控制复现", size=21)
    if VISUAL.exists():
        slide.shapes.add_picture(str(VISUAL), PptCm(1.0), PptCm(3.75), width=PptCm(23.4))
    add_textbox(slide, 1.0, 12.5, 23.4, 0.6, "目标：理解训练流程、部署方式、参数设置、结果判断和各脚本作用", size=16)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "1. 强化学习本质")
    add_pill(slide, 1.0, 2.2, 3.2, 1.0, "智能体")
    add_pill(slide, 5.0, 2.2, 3.2, 1.0, "环境")
    add_pill(slide, 9.0, 2.2, 3.2, 1.0, "动作")
    add_pill(slide, 13.0, 2.2, 3.2, 1.0, "奖励")
    add_pill(slide, 17.0, 2.2, 3.2, 1.0, "策略更新")
    add_bullet_box(slide, 1.0, 4.0, 22.5, 5.5, [
        "强化学习不是直接写死控制公式，而是在环境中反复试错。",
        "策略根据当前状态输出动作，环境推进一步并给出奖励。",
        "PPO 根据大量状态-动作-奖励样本更新神经网络。",
        "训练目标是让策略逐渐学会让无人机靠近并稳定在目标高度。",
    ], size=21)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "2. 本项目训练任务")
    add_bullet_box(slide, 1.0, 2.1, 11.3, 4.2, [
        "任务：单架四旋翼无人机悬停控制。",
        "目标位置：x = 0, y = 0, z = 1。",
        "环境：HoverAviary。",
        "算法：Stable-Baselines3 PPO。",
    ], size=20)
    slide.shapes.add_picture(str(target_formula), PptCm(14.0), PptCm(2.3), width=PptCm(7.0))
    slide.shapes.add_picture(str(reward_formula), PptCm(12.4), PptCm(5.0), width=PptCm(10.0))
    add_textbox(slide, 1.0, 8.4, 22.8, 2.0, "奖励函数的含义：无人机越接近目标点，奖励越高；偏离越大，奖励越低。", size=22, bold=True, fill=LIGHT)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "3. 训练工作流")
    slide.shapes.add_picture(str(workflow), PptCm(0.8), PptCm(2.0), width=PptCm(23.8))
    add_bullet_box(slide, 1.2, 8.2, 22.5, 3.2, [
        "reset 创建初始场景，step 执行动作并推进物理仿真。",
        "每一步返回 obs、reward、terminated、truncated。",
        "PPO 收集一批经验后更新策略，并保存评估表现更好的模型。",
    ], size=19)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "4. 软件与库的分工")
    slide.shapes.add_picture(str(roles), PptCm(1.1), PptCm(2.0), width=PptCm(22.5))

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "5. 关键脚本分工")
    items = [
        "reproducibility/docker/Dockerfile.repro：搭建稳定运行环境。",
        "HoverAviary.py：定义目标位置、奖励和结束条件。",
        "BaseRLAviary.py：定义观测/动作空间和动作预处理。",
        "experiments/hover_rl_reproduction/scripts/reproduce_hover_short.py：最小复现训练入口。",
        "learn.py：官方长训练、评估和 best_model 保存。",
        "play.py / render_policy_scene.py：策略播放、截图和录像。",
    ]
    add_bullet_box(slide, 1.1, 2.0, 22.5, 8.7, items, size=21)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "6. 部署环境")
    add_bullet_box(slide, 1.0, 2.1, 22.8, 3.0, [
        "推荐 Docker：避免 Windows 原生 Python、PyBullet 编译和依赖版本问题。",
        "镜像内环境：Linux + Python 3.10 + PyBullet + Gymnasium + Stable-Baselines3。",
    ], size=20)
    add_textbox(slide, 1.0, 5.5, 22.8, 2.8, "cd E:\\1-AI辅助工作\\科研项目\\强化学习\\gym-pybullet-drones\n\ndocker build -f reproducibility/docker/Dockerfile.repro -t gym-pybullet-drones-repro .", size=17, fill=LIGHT)
    add_textbox(slide, 1.0, 9.1, 22.8, 1.5, "部署完成后，所有训练输出都保存回本机 results 目录，便于后续分析和展示。", size=21, bold=True)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "7. 参数设置与启动训练")
    add_bullet_box(slide, 1.0, 2.1, 11.5, 4.0, [
        "OBS_TYPE = KIN：使用运动学状态。",
        "ACT_TYPE = ONE_D_RPM：一维动作控制四个电机共同变化。",
        "--timesteps：训练步数。",
        "--eval-episodes：评估回合数。",
        "--rollout-steps：导出轨迹长度。",
    ], size=18)
    add_textbox(slide, 13.0, 2.1, 10.5, 4.7, "docker run --rm ...\npython experiments/hover_rl_reproduction/scripts/reproduce_hover_short.py\n  --timesteps 2048\n  --eval-episodes 3\n  --rollout-steps 240\n  --output-folder experiments/hover_rl_reproduction/results/repro_hover_short", size=16, fill=LIGHT)
    add_textbox(slide, 1.0, 8.0, 22.5, 1.6, "2048 步用于验证链路；正式实验建议提高到 1e6 以上，并使用多个随机种子。", size=22, bold=True)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "8. 如何判断训练结果")
    bullets = [
        "mean_reward 越高通常越好，官方单机 one_d_rpm 目标奖励约 474。",
        "final_z 应接近目标高度 1.0。",
        "rollout.csv 中 x/y 应接近 0，z 应接近 1。",
        "不应频繁 truncated，动作不应剧烈震荡。",
        "可视化视频中应看到无人机接近目标高度并保持稳定。",
    ]
    add_bullet_box(slide, 1.0, 2.0, 13.0, 6.8, bullets, size=19)
    if summary:
        stats = (
            f"当前短训练结果\n"
            f"timesteps: {summary.get('timesteps')}\n"
            f"mean_reward: {float(summary.get('mean_reward')):.4f}\n"
            f"std_reward: {float(summary.get('std_reward')):.4f}\n"
            f"final_z: {float(summary.get('final_z')):.4f}\n"
            f"target_z: 1.0000"
        )
        add_textbox(slide, 15.0, 2.0, 8.2, 5.2, stats, size=19, fill=LIGHT)
    add_textbox(slide, 1.0, 10.3, 22.6, 1.4, "结论：当前短训练证明流程跑通，但还不是高质量控制器。", size=23, bold=True)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "9. 可视化展示")
    if VISUAL.exists():
        slide.shapes.add_picture(str(VISUAL), PptCm(1.0), PptCm(2.0), width=PptCm(14.5))
    add_bullet_box(slide, 16.0, 2.0, 7.5, 5.8, [
        "PyBullet 提供 3D 仿真场景。",
        "render_policy_scene.py 生成截图、GIF 和 MP4。",
        "适合汇报展示：不依赖现场打开 GUI。",
        "短训练视频展示的是链路验证，不是最终性能。",
    ], size=18)
    add_textbox(slide, 1.0, 10.6, 22.8, 1.0, "素材路径：docs/visualizations/hover_policy_scene_short/policy_scene_video.mp4", size=18, fill=LIGHT)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "10. 后续正式研究路线")
    add_bullet_box(slide, 1.0, 2.0, 22.8, 8.8, [
        "把训练步数从 2048 扩展到 100000、1000000 甚至更高。",
        "使用多个随机种子重复训练，验证稳定性。",
        "与 PID / DSLPIDControl 做同条件对比。",
        "优化奖励函数：加入姿态稳定、能耗和动作平滑惩罚。",
        "加入风扰、噪声、电机延迟，测试鲁棒性。",
        "进一步考虑 SITL、ROS2、PX4 或实机部署。",
    ], size=21)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "11. 一句话总结")
    add_textbox(slide, 1.4, 3.0, 22.5, 3.8, "本项目用 PyBullet 搭建四旋翼仿真环境，用 Gymnasium 封装强化学习接口，用 Stable-Baselines3 的 PPO 算法训练无人机悬停控制策略。", size=30, bold=True, fill=LIGHT)
    add_textbox(slide, 1.4, 8.0, 22.5, 2.5, "训练结果通过模型文件、summary.json、rollout.csv、曲线和 3D 场景视频进行解释和展示。", size=24)

    prs.save(PPT_OUT)


if __name__ == "__main__":
    build_word()
    build_ppt()
    print(WORD_OUT)
    print(PPT_OUT)
